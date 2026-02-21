# backend/services/ingest.py
from __future__ import annotations

import logging
import os
from pathlib import Path
from typing import List

from sqlalchemy.orm import Session

from backend.models import Document, DocumentChunk
from backend.services.embeddings import embed_texts
from backend.services.vector_index import vector_index

from langchain_core.documents import Document as LangchainDocument
from typing import Optional, Dict

logger = logging.getLogger(__name__)

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None  # type: ignore

try:
    import openpyxl # openpyxl
except ImportError:
    openpyxl = None  # type: ignore

try:
    from pptx import Presentation  # python-pptx
except ImportError:
    Presentation = None  # type: ignore

try:
    from docx2pdf import convert
except ImportError:
    docx2pdf = None  # type: ignore



def _read_pdf_pages(path: Path) -> List[tuple[int, str]]:
    """
    Returns list of (page_number, text) using PyMuPDF.
    """
    if fitz is None:
        raise RuntimeError("PyMuPDF (fitz) is not installed")

    doc = fitz.open(path)
    pages: List[tuple[int, str]] = []
    for i, page in enumerate(doc):
        text = page.get_text("text")
        text = text.strip()
        if text:
            pages.append((i + 1, text))
    doc.close()
    return pages


def _read_txt_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")



def _chunk_long_text(text: str, max_chars: int = 1000) -> List[str]:
    """
    Simple text chunker: split on paragraphs/newlines and group into ~max_chars chunks.
    Used for TXT and PDF page text.
    """
    pieces: List[str] = []
    buf: List[str] = []
    current_len = 0

    for para in text.splitlines():
        para = para.strip()
        if not para:
            continue
        if current_len + len(para) + 1 > max_chars and buf:
            pieces.append("\n".join(buf))
            buf = [para]
            current_len = len(para)
        else:
            buf.append(para)
            current_len += len(para) + 1

    if buf:
        pieces.append("\n".join(buf))

    return pieces


def _read_pptx_slides(path: Path) -> List[tuple[int, str]]:
    """
    Read PPTX slides as list of (slide_number, text) tuples.
    Each slide's text is combined from all text shapes.
    """
    if Presentation is None:
        raise RuntimeError("python-pptx is not installed")
    
    slides_with_text: List[tuple[int, str]] = []
    presentation = Presentation(path)
    
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    text_run = "".join(run.text for run in paragraph.runs)
                    if text_run.strip():
                        slide_text_parts.append(text_run.strip())
            elif hasattr(shape, "text"):
                if shape.text.strip():
                    slide_text_parts.append(shape.text.strip())
        
        if slide_text_parts:
            slide_text = "\n".join(slide_text_parts)
            slides_with_text.append((slide_number, slide_text))
    
    return slides_with_text

def ingest_document(db: Session, document: Document) -> None:
    """
    Read the document's file from disk, extract text, create DocumentChunk rows.

    - PDF: per page, further split into ~1200-char chunks.
    - DOCX/DOC: paragraph-aware chunking (~800 chars) with heading detection.
    - TXT: simple ~1000-char chunks.

    NOTE: This still only works for text-based PDFs, not scanned-image PDFs.
    OCR will be added separately.
    """
    path = Path(document.file_path)
    ext = document.file_extension.lower()

    logger.info("Ingesting document %s (%s)", document.id, path)

    db.query(DocumentChunk).filter(DocumentChunk.document_id == document.id).delete()

    chunks: List[DocumentChunk] = []
    chunk_index = 0

    try:
        if ext == "pdf":
            if not path.exists():
                raise FileNotFoundError(path)
            chunk_index = pdf_chunking(document, path, chunks, chunk_index)

        elif ext in ("docx", "doc"):
            if not path.exists():
                raise FileNotFoundError(path)
            try:
                convert(str(path), str(path.with_suffix(".pdf")))
                chunk_index = pdf_chunking(document, path.with_suffix(".pdf"), chunks, chunk_index)
            except Exception as e:
                logger.warning("DOCX to PDF conversion failed for document %s: %s", document.id, e)
            finally:
                # 3. Clean up the physical file
                pdf_path = path.with_suffix(".pdf")
                if os.path.exists(pdf_path):
                    os.remove(pdf_path) # Deletes the file from the system
        
        elif ext in ("xlsx", "xls"):
            if openpyxl is None:
                raise RuntimeError("openpyxl is not installed, cannot ingest Excel files")
            if not path.exists():
                raise FileNotFoundError(path)
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            for sheet_number, sheet in enumerate(wb.worksheets, start=1):
                sheet_text = []
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                    sheet_text.append(row_text)
                full_text = "\n".join(sheet_text)
                for piece in _chunk_long_text(full_text, max_chars=1000):
                    if not piece.strip():
                        continue
                    chunk = DocumentChunk(
                        document_id=document.id,
                        chunk_index=chunk_index,
                        text=piece,
                        page_start=sheet_number,
                        page_end=sheet_number,
                    )
                    chunks.append(chunk)
                    chunk_index += 1
        
        elif ext in ("pptx", "ppt"):         
            if not path.exists():
                raise FileNotFoundError(path)
            slides_with_text = _read_pptx_slides(path)
            for slide_no, slide_text in slides_with_text:
                for piece in _chunk_long_text(slide_text, max_chars=1000):
                    if not piece.strip():
                        continue
                    chunk = DocumentChunk(
                        document_id=document.id,
                        chunk_index=chunk_index,
                        text=piece,
                        page_start=slide_no,
                        page_end=slide_no,
                    )
                    chunks.append(chunk)
                    chunk_index += 1

        elif ext in ("txt",):
            if not path.exists():
                raise FileNotFoundError(path)
            full_text = _read_txt_text(path)
            for piece in _chunk_long_text(full_text, max_chars=1000):
                if not piece.strip():
                    continue
                chunk = DocumentChunk(
                    document_id=document.id,
                    chunk_index=chunk_index,
                    text=piece,
                    page_start=None,
                    page_end=None,
                )
                chunks.append(chunk)
                chunk_index += 1

        else:
            logger.warning("Unsupported extension %s for document %s", ext, document.id)
            return

        for c in chunks:
            db.add(c)

        db.commit()
        logger.info("Ingested %d chunks for document %s", len(chunks), document.id)

        try:
            chunk_ids = [c.id for c in chunks]
            vector_index.add_chunks(db, chunk_ids, embed_texts)
            logger.info(
                "Added %d chunks to vector index for document %s",
                len(chunk_ids),
                document.id,
            )
        except Exception as e:
            logger.exception(
                "Error adding chunks to vector index for document %s: %s",
                document.id,
                e,
            )

    except Exception as e:
        db.rollback()
        logger.exception("Error ingesting document %s: %s", document.id, e)

##########################################################################
# Alternative ingestion using LangChain 
##########################################################################

def parsing_loader(file_path: str, file_extension: str)->Optional[List[LangchainDocument]]:
    # from langchain_community.document_loaders import Docx2txtLoader
    """
    Parse a file using the appropriate LangChain document loader.

    Supported file extensions include docx, doc, pdf, txt, xlsx, xls, csv, ppt, and pptx.

    :param file_path: The path to the file to be parsed
    :param file_extension: The extension of the file to be parsed
    :return: A list of Langchain documents parsed from the file, or None if the extension is unsupported
    """
    from langchain_community.document_loaders import PyMuPDFLoader
    from langchain_community.document_loaders import TextLoader
    from langchain_community.document_loaders.base import BaseLoader
    from langchain_community.document_loaders import (
        UnstructuredExcelLoader, UnstructuredCSVLoader, UnstructuredPowerPointLoader
    )
    
    parsor_collections: Dict[str, BaseLoader] = {
        'docx': PyMuPDFLoader,
        'doc': PyMuPDFLoader,
        'pdf': PyMuPDFLoader,
        'txt': TextLoader,
        'xlsx': UnstructuredExcelLoader,  
        'xls': UnstructuredExcelLoader,
        'csv': UnstructuredCSVLoader,
        'pptx': UnstructuredPowerPointLoader,
        'ppt': UnstructuredPowerPointLoader
    }

    parsor:Optional[BaseLoader] = parsor_collections.get(file_extension.lower())

    if parsor:

        kwargs = dict(file_path=file_path)

        if file_extension == 'txt': kwargs['encoding'] = 'utf-8'
        if file_extension in ('xlsx', 'xls', 'pptx', 'ppt'): kwargs['mode'] = 'elements'

        loader:BaseLoader = parsor(**kwargs)

        documents = loader.load()
        logging.info("Parsed %d pages/chunks from document %s", len(documents), file_path)
        return documents
    
    else: logging.warning("Unsupported extension %s for document %s", file_extension, file_path)
    return None


def chunking_strategy(document: List[LangchainDocument])->List[LangchainDocument]:
    """
    Split a list of Langchain documents into a new list of Langchain documents
    where each document is split into chunks of approximately 1000 characters
    with an overlap of 200 characters between adjacent chunks.

    :param document: A list of Langchain documents to be split
    :return: A list of Langchain documents, each split into chunks
    """
    from langchain_text_splitters import RecursiveCharacterTextSplitter

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        separators=["\n\n", "\n", " "],
    )

    return text_splitter.split_documents(
        document
    )


    
def ingest_document_langchain(db: Session, document: Document) -> None:
    """
    Read the document's file from disk, extract text, create DocumentChunk rows.

    - PDF: per page, further split into ~1200-char chunks.
    - DOCX/DOC: paragraph-aware chunking (~800 chars) with heading detection.
    - TXT: simple ~1000-char chunks.

    NOTE: This still only works for text-based PDFs, not scanned-image PDFs.
    OCR will be added separately.
    """

    path = Path(document.file_path)
    ext = document.file_extension.lower()

    logging.info("Ingesting document %s (%s)", document.id, path)

    db.query(DocumentChunk).filter(DocumentChunk.document_id == document.id).delete()

    chunks: List[DocumentChunk] = []
    chunk_index = 0

    try:
        documents = parsing_loader(str(path), ext)

        if documents:
            chunked_docs = chunking_strategy(documents)

            for c in chunked_docs:
                page_no = c.metadata.get("page") if c.metadata else None
                page_no = c.metadata.get("page_number") if page_no is None and c.metadata else page_no

                chunk = DocumentChunk(
                    document_id=document.id,
                    chunk_index=chunk_index,
                    text=c.page_content,
                    page_start=page_no,
                    page_end=page_no,
                )
                chunks.append(chunk)
                chunk_index += 1

        for c in chunks:
            db.add(c)

        db.commit()
        logging.info("Ingested %d chunks for document %s", len(chunks), document.id)

        try:
            chunk_ids = [c.id for c in chunks]
            vector_index.add_chunks(db, chunk_ids, embed_texts)
            logging.info(
                "Added %d chunks to vector index for document %s",
                len(chunk_ids),
                document.id,
            )
        except Exception as e:
            logging.exception(
                "Error adding chunks to vector index for document %s: %s",
                document.id,
                e,
            )

    except Exception as e:
        db.rollback()
        logging.exception("Error ingesting document %s: %s", document.id, e)
